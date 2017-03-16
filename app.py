#wikipedia api is easy but not for advanced usage, thus BeautifulSoup is used
from pptx import Presentation
prs = Presentation()
from pptx.util import Inches

import wikipedia
#wiki text scrap
from bs4 import BeautifulSoup
#html request
import requests
from cherrypy.lib import static
import cherrypy
import os
from jinja2 import Environment, FileSystemLoader

localDir = os.path.dirname(__file__)
absDir = os.path.join(os.getcwd(), localDir)

env = Environment(loader=FileSystemLoader('html'))

class Slide_1(object):
    def __init__(self,city_name):
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self.wiki_title(city_name)
        subtitle.text = self.wiki_subtitle(city_name)
        #img_path = self.wiki_image(city_name)
        left = Inches(4.5)
        top = Inches(1)
        #pic = slide.shapes.add_picture(img_path, left, top, width=1280160, height=1371600)

    def wiki_title(self,city_name):
        city = wikipedia.page(city_name)
        return city.title

    def wiki_subtitle(self,city_name):
        summary = wikipedia.summary(city_name, sentences=1)
        #shorten_summary = 
        #mit regex machen
        return summary

    def wiki_image(self,city_name):
        city = wikipedia.page(city_name)
        return city.images[0]

class Slide_2(object):
    def __init__(self, city_name):
        title_slide_layouts = prs.slide_layouts[1]
        slides = prs.slides.add_slide(title_slide_layouts)
        titles = slides.shapes.title
        subtitles = slides.placeholders[1]
        titles.text = "Hello 2"
        subtitles.text = "lala"

class Table_of_content(object):
    def __init__(self,city_name):
        #beautifulsoup library
        url= "https://en.wikipedia.org/wiki/" + city_name
        content = requests.get(url).content
        soup = BeautifulSoup(content, 'lxml')

    def all_content(self):
        tag = soup.find('div', {'class': 'toclevel1'})
        links=tag.findAll('a')
        content=[]

        for link in links:
            content.append(link.text)
        return content

class Contents(object):
    def __init__(self):
        title = self.title

class References(object):
    def __init__(self):
        title = self.title

class Landing_Page(object):
    @cherrypy.expose
    def index(self):
        tmpl = env.get_template('index.html')
        return tmpl.render()

    @cherrypy.expose
    def check_validity(self,city_name):
        try:
            return self.powerpoint(city_name)
        except wikipedia.exceptions.DisambiguationError as e:
            return e.options

    @cherrypy.expose
    def powerpoint(self, city_name):
        Slide_1(city_name)
        Slide_2(city_name)
        #Table_of_content(city_name)
        
        prs.save(city_name + '.pptx')
        #TODO catch exception 
        return self.download(city_name)

    @cherrypy.expose
    def download(self, city_name):
        path = os.path.join(absDir, city_name + ".pptx")
        result = static.serve_file(path, "application/x-download",
                                 "attachment", name = os.path.basename(path))
        #os.unlink(filename)
        return result


config = {
    'global': {
        'server.socket_host': '0.0.0.0',
        'server.socket_port': int(os.environ.get('PORT', 5000)),
    },
    '/assets': {
        'tools.staticdir.root': os.path.dirname(os.path.abspath(__file__)),
        'tools.staticdir.on': True,
        'tools.staticdir.dir': 'assets',
    }
}

cherrypy.quickstart(Landing_Page(), '/', config=config)
